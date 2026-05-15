from __future__ import annotations

import argparse
import json
import math
from datetime import datetime
from html import escape
from pathlib import Path
from typing import Iterable

import folium
import h3
import matplotlib
import numpy as np
import pandas as pd
import requests
from branca.colormap import LinearColormap
from branca.element import Element
from folium.plugins import MarkerCluster
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from sklearn.ensemble import RandomForestClassifier
from sklearn.metrics import (
    average_precision_score,
    balanced_accuracy_score,
    classification_report,
    confusion_matrix,
    precision_score,
    recall_score,
    roc_auc_score,
)
from sklearn.model_selection import train_test_split


matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402


ROOT = Path(__file__).resolve().parents[1]
DATA_PARQUET = ROOT / "data.parquet"
TARGET_PARQUET = ROOT / "target.parquet"

OUT = ROOT / "output"
DATA_OUT = OUT / "data"
CHART_OUT = OUT / "charts"
MAP_OUT = OUT / "maps"
PPTX_OUT = OUT / "pptx"
OSM_OUT = OUT / "osm"
OSM_CACHE = OSM_OUT / "osm_poi_raw.json"
OSM_META = OSM_OUT / "osm_fetch_meta.json"

VTB_BLUE = RGBColor(0, 63, 136)
TEXT = RGBColor(32, 42, 54)
MUTED = RGBColor(92, 103, 116)

RANDOM_STATE = 42

POI_CATEGORY_HELP = {
    "vtb_atm": (
        "vtb_atm",
        "Банкомат ВТБ",
        "Найденный в OSM банкомат ВТБ. Используется как признак текущего присутствия банка.",
        "vtb_presence_raw = vtb_atm + vtb_branch",
    ),
    "vtb_branch": (
        "vtb_branch",
        "Отделение ВТБ",
        "Найденное в OSM отделение/банк ВТБ. Зона скорее относится к усилению текущей сети.",
        "capacity_candidate = vtb_presence_raw > 0",
    ),
    "competitor_atm": (
        "competitor_atm",
        "Банкомат конкурента",
        "Прямой конкурентный ATM-сигнал. Подтверждает рыночный спрос рядом.",
        "competition_raw = competitor_atm + 0.7 * competitor_bank",
    ),
    "competitor_bank": (
        "competitor_bank",
        "Банк-конкурент",
        "Более мягкий сигнал финансовой активности, чем банкомат.",
        "competition_raw = competitor_atm + 0.7 * competitor_bank",
    ),
    "mall": (
        "mall",
        "Торговый центр",
        "Трафиковый POI. Усиливает потенциал зоны за счет регулярного клиентского потока.",
        "traffic_poi_raw = 2 * mall + 1.5 * transit + education",
    ),
    "transit": (
        "transit",
        "Транспортный объект",
        "Станция или вход метро/жд. Используется как прокси проходящего потока.",
        "traffic_poi_raw = 2 * mall + 1.5 * transit + education",
    ),
    "education": (
        "education",
        "Образовательный объект",
        "Вуз или колледж. Дает дополнительный регулярный дневной трафик.",
        "traffic_poi_raw = 2 * mall + 1.5 * transit + education",
    ),
}


def ensure_dirs() -> None:
    for path in [DATA_OUT, CHART_OUT, MAP_OUT, PPTX_OUT, OSM_OUT]:
        path.mkdir(parents=True, exist_ok=True)


def qnorm(values: Iterable[float]) -> pd.Series:
    s = pd.Series(values, dtype="float64")
    lo, hi = s.quantile(0.01), s.quantile(0.99)
    clipped = s.clip(lo, hi)
    if math.isclose(float(lo), float(hi)):
        return pd.Series(np.zeros(len(s)), index=s.index)
    return (clipped - lo) / (hi - lo)


def entropy_from_counts(values: pd.Series | np.ndarray) -> float:
    arr = np.asarray(values, dtype="float64")
    total = arr.sum()
    if total <= 0:
        return 0.0
    p = arr[arr > 0] / total
    return float(-(p * np.log(p)).sum())


def haversine_km(lat1: np.ndarray, lon1: np.ndarray, lat2: np.ndarray, lon2: np.ndarray) -> np.ndarray:
    radius = 6371.0088
    lat1_rad = np.radians(lat1)
    lon1_rad = np.radians(lon1)
    lat2_rad = np.radians(lat2)
    lon2_rad = np.radians(lon2)
    dlat = lat2_rad - lat1_rad
    dlon = lon2_rad - lon1_rad
    a = np.sin(dlat / 2) ** 2 + np.cos(lat1_rad) * np.cos(lat2_rad) * np.sin(dlon / 2) ** 2
    return 2 * radius * np.arcsin(np.sqrt(a))


def overpass_query(bbox: tuple[float, float, float, float]) -> str:
    south, west, north, east = bbox
    return f"""
[out:json][timeout:35];
(
  node["amenity"~"^(atm|bank)$"]({south},{west},{north},{east});
  way["amenity"~"^(atm|bank)$"]({south},{west},{north},{east});
  node["shop"="mall"]({south},{west},{north},{east});
  way["shop"="mall"]({south},{west},{north},{east});
  node["railway"~"^(station|subway_entrance)$"]({south},{west},{north},{east});
  way["railway"~"^(station|subway_entrance)$"]({south},{west},{north},{east});
  node["amenity"~"^(university|college)$"]({south},{west},{north},{east});
  way["amenity"~"^(university|college)$"]({south},{west},{north},{east});
);
out center tags;
"""


def fetch_osm(cell: pd.DataFrame, refresh: bool = False, no_osm: bool = False) -> tuple[dict[str, object] | None, dict[str, object]]:
    if no_osm:
        return None, {"skipped": True, "elements": 0, "timestamp_osm_base": None}
    if OSM_CACHE.exists() and not refresh:
        raw = json.loads(OSM_CACHE.read_text(encoding="utf-8"))
        meta = json.loads(OSM_META.read_text(encoding="utf-8")) if OSM_META.exists() else {"source": "cache"}
        meta["used_cache"] = True
        return raw, meta

    bbox = (
        float(cell["lat"].min() - 0.02),
        float(cell["lon"].min() - 0.02),
        float(cell["lat"].max() + 0.02),
        float(cell["lon"].max() + 0.02),
    )
    endpoints = [
        "https://overpass-api.de/api/interpreter",
        "https://overpass.kumi.systems/api/interpreter",
    ]
    errors = []
    for endpoint in endpoints:
        try:
            response = requests.post(
                endpoint,
                data={"data": overpass_query(bbox)},
                timeout=80,
                headers={"User-Agent": "vtb-atm-case-v2-ml/1.0"},
            )
            if response.ok:
                raw = response.json()
                OSM_CACHE.write_text(json.dumps(raw, ensure_ascii=False), encoding="utf-8")
                meta = {
                    "used_cache": False,
                    "endpoint": endpoint,
                    "bbox": bbox,
                    "elements": len(raw.get("elements", [])),
                    "timestamp_osm_base": raw.get("osm3s", {}).get("timestamp_osm_base"),
                }
                OSM_META.write_text(json.dumps(meta, ensure_ascii=False, indent=2), encoding="utf-8")
                return raw, meta
            errors.append({"endpoint": endpoint, "status": response.status_code, "text": response.text[:300]})
        except Exception as exc:  # noqa: BLE001
            errors.append({"endpoint": endpoint, "error": type(exc).__name__, "message": str(exc)})
    return None, {"used_cache": False, "errors": errors}


def classify_poi(tags: dict[str, object]) -> str | None:
    text = " ".join(str(tags.get(k, "")) for k in ["name", "brand", "operator", "network"]).lower()
    amenity = tags.get("amenity")
    shop = tags.get("shop")
    railway = tags.get("railway")
    if amenity == "atm":
        return "vtb_atm" if "втб" in text or "vtb" in text else "competitor_atm"
    if amenity == "bank":
        return "vtb_branch" if "втб" in text or "vtb" in text else "competitor_bank"
    if shop == "mall":
        return "mall"
    if railway in {"station", "subway_entrance"}:
        return "transit"
    if amenity in {"university", "college"}:
        return "education"
    return None


def osm_to_poi(raw: dict[str, object] | None, case_cells: set[str]) -> pd.DataFrame:
    rows: list[dict[str, object]] = []
    if raw is None:
        return pd.DataFrame(columns=["h3_index", "lat", "lon", "category", "name", "osm_id", "in_case_cells"])
    for element in raw.get("elements", []):
        tags = element.get("tags", {})
        lat = element.get("lat") or element.get("center", {}).get("lat")
        lon = element.get("lon") or element.get("center", {}).get("lon")
        category = classify_poi(tags)
        if lat is None or lon is None or category is None:
            continue
        h3_index = h3.latlng_to_cell(float(lat), float(lon), 9)
        rows.append(
            {
                "h3_index": h3_index,
                "lat": float(lat),
                "lon": float(lon),
                "category": category,
                "name": str(tags.get("name", "")),
                "osm_id": element.get("id"),
                "in_case_cells": h3_index in case_cells,
            }
        )
    return pd.DataFrame(rows)


def nearest_distance_km(cells: pd.DataFrame, poi: pd.DataFrame, category: str) -> pd.Series:
    subset = poi[(poi["category"] == category) & (poi["in_case_cells"])].copy()
    if subset.empty:
        return pd.Series(np.nan, index=cells.index)
    poi_lat = subset["lat"].to_numpy(dtype="float64")
    poi_lon = subset["lon"].to_numpy(dtype="float64")
    result: list[float] = []
    for _, row in cells.iterrows():
        distances = haversine_km(
            np.full(len(subset), row["lat"], dtype="float64"),
            np.full(len(subset), row["lon"], dtype="float64"),
            poi_lat,
            poi_lon,
        )
        result.append(float(distances.min()))
    return pd.Series(result, index=cells.index)


def count_within_km(cells: pd.DataFrame, poi: pd.DataFrame, categories: list[str], radius_km: float) -> pd.Series:
    subset = poi[(poi["category"].isin(categories)) & (poi["in_case_cells"])].copy()
    if subset.empty:
        return pd.Series(0, index=cells.index)
    poi_lat = subset["lat"].to_numpy(dtype="float64")
    poi_lon = subset["lon"].to_numpy(dtype="float64")
    counts: list[int] = []
    for _, row in cells.iterrows():
        distances = haversine_km(
            np.full(len(subset), row["lat"], dtype="float64"),
            np.full(len(subset), row["lon"], dtype="float64"),
            poi_lat,
            poi_lon,
        )
        counts.append(int((distances <= radius_km).sum()))
    return pd.Series(counts, index=cells.index)


def weighted_transaction_volatility(data: pd.DataFrame, totals: pd.DataFrame) -> pd.Series:
    merged = data.merge(totals[["h3_index", "avg_ticket_weighted"]], on="h3_index", how="left")
    n = merged["count"].astype(float)
    row_std = merged["std"].fillna(0).astype(float)
    row_mean = merged["avg"].astype(float)
    zone_mean = merged["avg_ticket_weighted"].astype(float)
    merged["variance_part"] = (n - 1).clip(lower=0) * (row_std**2) + n * ((row_mean - zone_mean) ** 2)
    numerator = merged.groupby("h3_index")["variance_part"].sum()
    denominator = totals.set_index("h3_index")["total_tx_count"].astype(float) - 1
    return np.sqrt((numerator / denominator.replace(0, np.nan)).fillna(0))


def load_data() -> tuple[pd.DataFrame, pd.DataFrame]:
    data = pd.read_parquet(DATA_PARQUET)
    target = pd.read_parquet(TARGET_PARQUET).drop_duplicates(["h3_index", "customer_id"])
    return data, target


def build_cell_features(data: pd.DataFrame, target: pd.DataFrame, poi: pd.DataFrame) -> tuple[pd.DataFrame, dict[str, object]]:
    base = (
        data.groupby("h3_index")
        .agg(
            row_count=("customer_id", "size"),
            unique_customers=("customer_id", "nunique"),
            total_tx_count=("count", "sum"),
            total_tx_sum=("sum", "sum"),
            avg_ticket_mean=("avg", "mean"),
            max_transaction=("max", "max"),
            count_distinct=("count_distinct", "sum"),
            unique_mcc_count=("mcc_code", "nunique"),
            active_time_buckets=("datetime_id", "nunique"),
        )
        .reset_index()
    )
    base["tx_per_customer"] = base["total_tx_count"] / base["unique_customers"].replace(0, np.nan)
    base["sum_per_customer"] = base["total_tx_sum"] / base["unique_customers"].replace(0, np.nan)
    base["avg_ticket_weighted"] = base["total_tx_sum"] / base["total_tx_count"].replace(0, np.nan)
    base["lat"] = base["h3_index"].map(lambda idx: h3.cell_to_latlng(idx)[0])
    base["lon"] = base["h3_index"].map(lambda idx: h3.cell_to_latlng(idx)[1])

    volatility = weighted_transaction_volatility(data, base).rename("transaction_volatility").reset_index()
    base = base.merge(volatility, on="h3_index", how="left")

    avg_q25 = data["avg"].quantile(0.25)
    avg_q75 = data["avg"].quantile(0.75)
    small = data[data["avg"] <= avg_q25].groupby("h3_index")["count"].sum().rename("small_ticket_count")
    large = data[data["avg"] >= avg_q75].groupby("h3_index")["count"].sum().rename("large_ticket_count")
    base = base.merge(small.reset_index(), on="h3_index", how="left").merge(large.reset_index(), on="h3_index", how="left")
    base[["small_ticket_count", "large_ticket_count"]] = base[["small_ticket_count", "large_ticket_count"]].fillna(0)
    base["small_ticket_share"] = base["small_ticket_count"] / base["total_tx_count"].replace(0, np.nan)
    base["large_ticket_share"] = base["large_ticket_count"] / base["total_tx_count"].replace(0, np.nan)

    mcc_counts = data.groupby(["h3_index", "mcc_code"])["count"].sum().reset_index()
    mcc_entropy = mcc_counts.groupby("h3_index")["count"].apply(entropy_from_counts).rename("mcc_entropy").reset_index()
    base = base.merge(mcc_entropy, on="h3_index", how="left")
    top_mcc = mcc_counts.groupby("h3_index")["count"].max().rename("top_mcc_count").reset_index()
    base = base.merge(top_mcc, on="h3_index", how="left")
    base["top_mcc_share"] = base["top_mcc_count"] / base["total_tx_count"].replace(0, np.nan)

    all_mcc_codes = sorted(int(code) for code in data["mcc_code"].dropna().unique())
    mcc_pivot = mcc_counts.pivot(index="h3_index", columns="mcc_code", values="count").fillna(0)
    for code in all_mcc_codes:
        if code not in mcc_pivot.columns:
            mcc_pivot[code] = 0.0
    mcc_share = mcc_pivot[all_mcc_codes].div(mcc_pivot[all_mcc_codes].sum(axis=1).replace(0, np.nan), axis=0).fillna(0)
    mcc_share.columns = [f"mcc_vector_{code}" for code in all_mcc_codes]
    base = base.merge(mcc_share.reset_index(), on="h3_index", how="left")

    time_counts = data.groupby(["h3_index", "datetime_id"])["count"].sum().reset_index()
    time_pivot = time_counts.pivot(index="h3_index", columns="datetime_id", values="count").fillna(0)
    time_entropy = (time_pivot.apply(entropy_from_counts, axis=1) / math.log(max(2, time_pivot.shape[1]))).rename("time_entropy")
    peak_count = time_pivot.max(axis=1).rename("peak_bucket_count")
    base = base.merge(time_entropy.reset_index(), on="h3_index", how="left").merge(peak_count.reset_index(), on="h3_index", how="left")
    base["peak_bucket_share"] = base["peak_bucket_count"] / base["total_tx_count"].replace(0, np.nan)
    base["off_peak_activity"] = 1 - base["peak_bucket_share"]

    target_counts = target.groupby("h3_index")["customer_id"].nunique().rename("atm_active_customers")
    base = base.merge(target_counts.reset_index(), on="h3_index", how="left")
    base["atm_active_customers"] = base["atm_active_customers"].fillna(0).astype(int)
    base["cell_y"] = (base["atm_active_customers"] > 0).astype(int)
    base["atm_penetration_observed"] = base["atm_active_customers"] / base["unique_customers"].replace(0, np.nan)

    strong_threshold = base["atm_active_customers"].quantile(0.8)
    strong_cells = set(base.loc[base["atm_active_customers"] >= strong_threshold, "h3_index"])
    strong_counts = mcc_counts[mcc_counts["h3_index"].isin(strong_cells)].groupby("mcc_code")["count"].sum()
    global_counts = mcc_counts.groupby("mcc_code")["count"].sum()
    lift = (strong_counts / strong_counts.sum()) / (global_counts / global_counts.sum())
    lift_norm = qnorm(lift.reindex(all_mcc_codes).replace([np.inf, -np.inf], np.nan).fillna(0)).to_numpy()
    base["atm_affinity_mcc_score"] = mcc_share[[f"mcc_vector_{code}" for code in all_mcc_codes]].to_numpy().dot(lift_norm)

    categories = ["vtb_atm", "vtb_branch", "competitor_atm", "competitor_bank", "mall", "transit", "education"]
    for category in categories:
        counts = poi[(poi["category"] == category) & (poi["in_case_cells"])].groupby("h3_index").size()
        base[category] = base["h3_index"].map(counts).fillna(0).astype(int)

    cells = base[["h3_index", "lat", "lon"]].copy()
    base["distance_to_nearest_vtb_atm"] = nearest_distance_km(cells, poi, "vtb_atm")
    base["distance_to_nearest_competitor_atm"] = nearest_distance_km(cells, poi, "competitor_atm")
    base["competitor_atm_count_500m"] = count_within_km(cells, poi, ["competitor_atm"], 0.5)
    base["poi_count_500m"] = count_within_km(cells, poi, categories, 0.5)
    base["metro_distance"] = nearest_distance_km(cells, poi, "transit")
    base["mall_count_500m"] = count_within_km(cells, poi, ["mall"], 0.5)
    base["university_count_500m"] = count_within_km(cells, poi, ["education"], 0.5)

    base["traffic_poi_raw"] = 2 * base["mall"] + 1.5 * base["transit"] + base["education"]
    base["competition_raw"] = base["competitor_atm"] + 0.7 * base["competitor_bank"]
    base["vtb_presence_raw"] = base["vtb_atm"] + base["vtb_branch"]
    base["coverage_gap"] = (base["vtb_presence_raw"] == 0).astype(int)

    diagnostics = {
        "data_rows": int(len(data)),
        "target_rows_raw": int(pd.read_parquet(TARGET_PARQUET).shape[0]),
        "target_rows_unique": int(len(target)),
        "h3_cells_data": int(base["h3_index"].nunique()),
        "h3_cells_target": int(target["h3_index"].nunique()),
        "h3_positive_cells_in_data": int(base["cell_y"].sum()),
        "h3_negative_cells_in_data": int((base["cell_y"] == 0).sum()),
        "data_unique_pairs": int(data[["h3_index", "customer_id"]].drop_duplicates().shape[0]),
        "target_pairs_in_data": int(
            data[["h3_index", "customer_id"]]
            .drop_duplicates()
            .merge(target.assign(y=1), on=["h3_index", "customer_id"], how="inner")
            .shape[0]
        ),
        "target_pairs_absent_in_data": int(
            target.merge(
                data[["h3_index", "customer_id"]].drop_duplicates().assign(in_data=1),
                on=["h3_index", "customer_id"],
                how="left",
            )["in_data"]
            .isna()
            .sum()
        ),
        "h3_resolution": sorted({int(h3.get_resolution(idx)) for idx in base["h3_index"]}),
        "lat_range": [float(base["lat"].min()), float(base["lat"].max())],
        "lon_range": [float(base["lon"].min()), float(base["lon"].max())],
        "mcc_codes": all_mcc_codes,
        "datetime_buckets": sorted(int(x) for x in data["datetime_id"].dropna().unique()),
    }
    return base, diagnostics


def model_features(cell: pd.DataFrame) -> list[str]:
    excluded = {
        "h3_index",
        "cell_y",
        "atm_active_customers",
        "atm_penetration_observed",
        "atm_affinity_mcc_score",
        "small_ticket_count",
        "large_ticket_count",
        "top_mcc_count",
        "peak_bucket_count",
    }
    return [col for col in cell.columns if col not in excluded and pd.api.types.is_numeric_dtype(cell[col])]


def train_model(cell: pd.DataFrame) -> tuple[RandomForestClassifier, pd.DataFrame, dict[str, object], pd.DataFrame]:
    feature_cols = model_features(cell)
    x = cell[feature_cols].copy()
    y = cell["cell_y"].astype(int)

    fill_values = {}
    for col in feature_cols:
        if x[col].isna().any():
            if col.startswith("distance_") or col.endswith("_distance"):
                max_value = x[col].max(skipna=True)
                fill_values[col] = 999.0 if pd.isna(max_value) else float(max_value + 1)
            else:
                median_value = x[col].median(skipna=True)
                fill_values[col] = 0.0 if pd.isna(median_value) else float(median_value)
    x = x.fillna(fill_values)

    x_train, x_test, y_train, y_test = train_test_split(
        x, y, test_size=0.25, random_state=RANDOM_STATE, stratify=y
    )
    model = RandomForestClassifier(
        n_estimators=500,
        max_depth=12,
        min_samples_leaf=4,
        class_weight="balanced_subsample",
        random_state=RANDOM_STATE,
        n_jobs=-1,
    )
    model.fit(x_train, y_train)
    proba_test = model.predict_proba(x_test)[:, 1]
    pred_test = (proba_test >= 0.5).astype(int)
    top20_cut = np.quantile(proba_test, 0.8)
    pred_top20 = (proba_test >= top20_cut).astype(int)
    metrics = {
        "model": "RandomForestClassifier",
        "target": "cell_y = 1 if h3_index has at least one target customer",
        "features": feature_cols,
        "train_rows": int(len(x_train)),
        "test_rows": int(len(x_test)),
        "positive_rate_all": float(y.mean()),
        "roc_auc": float(roc_auc_score(y_test, proba_test)),
        "average_precision": float(average_precision_score(y_test, proba_test)),
        "balanced_accuracy_at_0_5": float(balanced_accuracy_score(y_test, pred_test)),
        "precision_at_0_5": float(precision_score(y_test, pred_test, zero_division=0)),
        "recall_at_0_5": float(recall_score(y_test, pred_test, zero_division=0)),
        "precision_top20pct": float(precision_score(y_test, pred_top20, zero_division=0)),
        "recall_top20pct": float(recall_score(y_test, pred_top20, zero_division=0)),
        "confusion_matrix_at_0_5": confusion_matrix(y_test, pred_test).tolist(),
        "classification_report_at_0_5": classification_report(y_test, pred_test, output_dict=True, zero_division=0),
        "fill_values": fill_values,
    }
    all_proba = model.predict_proba(x)[:, 1]
    scored = cell.copy()
    scored["ml_atm_probability"] = all_proba
    scored["ml_probability_rank"] = scored["ml_atm_probability"].rank(ascending=False, method="first").astype(int)

    importances = (
        pd.DataFrame({"feature": feature_cols, "importance": model.feature_importances_})
        .sort_values("importance", ascending=False)
        .reset_index(drop=True)
    )
    return model, scored, metrics, importances


def add_business_scores(scored: pd.DataFrame) -> pd.DataFrame:
    result = scored.copy()
    result["n_customers"] = qnorm(np.log1p(result["unique_customers"])).values
    result["n_tx"] = qnorm(np.log1p(result["total_tx_count"])).values
    result["n_sum"] = qnorm(np.log1p(result["total_tx_sum"])).values
    result["n_mcc"] = qnorm(result["mcc_entropy"] / np.log(result["unique_mcc_count"].clip(lower=2))).fillna(0).values
    result["n_time"] = qnorm(0.7 * result["time_entropy"] + 0.3 * qnorm(result["active_time_buckets"]).values).fillna(0).values
    result["demand_score_v2"] = (
        0.34 * result["n_customers"]
        + 0.27 * result["n_tx"]
        + 0.19 * result["n_sum"]
        + 0.12 * result["n_mcc"]
        + 0.08 * result["n_time"]
    )
    result["n_traffic"] = qnorm(np.log1p(result["traffic_poi_raw"])).values
    result["n_competition"] = qnorm(np.log1p(result["competition_raw"])).values
    distance_fill = result["distance_to_nearest_vtb_atm"].max(skipna=True)
    if pd.isna(distance_fill):
        distance_fill = 999.0
    result["n_distance_to_vtb"] = qnorm(np.log1p(result["distance_to_nearest_vtb_atm"].fillna(distance_fill))).values
    result["coverage_gap_score"] = (
        0.35 * result["coverage_gap"]
        + 0.25 * result["n_distance_to_vtb"]
        + 0.20 * result["n_traffic"]
        + 0.20 * result["n_competition"]
    )
    result["placement_score"] = 100 * (
        0.55 * result["ml_atm_probability"]
        + 0.25 * result["demand_score_v2"]
        + 0.20 * result["coverage_gap_score"]
    )
    anomaly = (result["unique_customers"] < 10) & (result["total_tx_sum"] > result["total_tx_sum"].quantile(0.95))
    result["recommendation_type"] = "Пилотная новая точка"
    result.loc[(result["cell_y"] == 0) & (result["ml_atm_probability"] >= result["ml_atm_probability"].quantile(0.9)), "recommendation_type"] = (
        "Латентный спрос: модель выше target"
    )
    result.loc[(result["cell_y"] == 1) & (result["vtb_presence_raw"] == 0), "recommendation_type"] = "Подтвержденный ATM-спрос без ВТБ"
    result.loc[(result["competition_raw"] >= 2) & (result["vtb_presence_raw"] == 0), "recommendation_type"] = "Конкурентный кластер без ВТБ"
    result.loc[result["vtb_presence_raw"] > 0, "recommendation_type"] = "Усиление существующей сети"
    result.loc[anomaly, "recommendation_type"] = "Проверить аномалию высокого чека"
    result["new_atm_candidate"] = (result["vtb_presence_raw"] == 0) & (~anomaly)
    result["capacity_candidate"] = result["vtb_presence_raw"] > 0
    result["priority_rank"] = result["placement_score"].rank(ascending=False, method="first").astype(int)
    result["why_recommended"] = result.apply(explain_row, axis=1)
    return result.sort_values("priority_rank")


def explain_row(row: pd.Series) -> str:
    reasons = []
    if row["ml_atm_probability"] >= 0.75:
        reasons.append("высокая ML-вероятность ATM-сигнала")
    elif row["ml_atm_probability"] >= 0.5:
        reasons.append("заметная ML-вероятность ATM-сигнала")
    if row["unique_customers"] >= 500:
        reasons.append("широкая клиентская база")
    elif row["unique_customers"] >= 150:
        reasons.append("заметная клиентская база")
    if row["total_tx_count"] >= 2000:
        reasons.append("много операций")
    if row["total_tx_sum"] >= 2_000_000:
        reasons.append("крупный оборот")
    if row["unique_mcc_count"] >= 10:
        reasons.append("разнообразные MCC-сценарии")
    if row["cell_y"] == 1:
        reasons.append("есть ATM-сигнал в target")
    if row["competition_raw"] >= 2:
        reasons.append("спрос подтвержден конкурентами")
    if row["coverage_gap"] == 1:
        reasons.append("нет найденного ВТБ-присутствия в OSM")
    return "; ".join(reasons[:6]) or "сбалансированный потенциал зоны"


def export_artifacts(scored: pd.DataFrame, poi: pd.DataFrame, diagnostics: dict[str, object], model_metrics: dict[str, object], importances: pd.DataFrame, osm_meta: dict[str, object]) -> dict[str, Path]:
    main_cols = [
        "priority_rank",
        "h3_index",
        "lat",
        "lon",
        "placement_score",
        "ml_atm_probability",
        "ml_probability_rank",
        "cell_y",
        "atm_active_customers",
        "demand_score_v2",
        "coverage_gap_score",
        "recommendation_type",
        "why_recommended",
        "new_atm_candidate",
        "capacity_candidate",
        "unique_customers",
        "total_tx_count",
        "total_tx_sum",
        "tx_per_customer",
        "sum_per_customer",
        "avg_ticket_weighted",
        "transaction_volatility",
        "unique_mcc_count",
        "mcc_entropy",
        "top_mcc_share",
        "active_time_buckets",
        "time_entropy",
        "peak_bucket_share",
        "vtb_atm",
        "vtb_branch",
        "competitor_atm",
        "competitor_bank",
        "mall",
        "transit",
        "education",
        "distance_to_nearest_vtb_atm",
        "distance_to_nearest_competitor_atm",
    ]
    paths = {
        "zone_scores": DATA_OUT / "v2_zone_scores_ml.csv",
        "shortlist": DATA_OUT / "v2_shortlist_new_atm.csv",
        "capacity": DATA_OUT / "v2_shortlist_capacity_existing_network.csv",
        "importances": DATA_OUT / "v2_model_feature_importance.csv",
        "poi": DATA_OUT / "v2_osm_poi_by_h3.csv",
        "diagnostics": DATA_OUT / "v2_diagnostics.json",
        "full_features": DATA_OUT / "v2_cell_features_full.csv",
    }
    scored.to_csv(paths["full_features"], index=False, encoding="utf-8-sig")
    scored[main_cols].to_csv(paths["zone_scores"], index=False, encoding="utf-8-sig")
    scored[scored["new_atm_candidate"]].head(30)[main_cols].to_csv(paths["shortlist"], index=False, encoding="utf-8-sig")
    scored[scored["capacity_candidate"]].head(20)[main_cols].to_csv(paths["capacity"], index=False, encoding="utf-8-sig")
    importances.to_csv(paths["importances"], index=False, encoding="utf-8-sig")
    poi.to_csv(paths["poi"], index=False, encoding="utf-8-sig")
    payload = {"case_data": diagnostics, "osm": osm_meta, "model_metrics": model_metrics}
    paths["diagnostics"].write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return paths


def save_charts(scored: pd.DataFrame, importances: pd.DataFrame, model_metrics: dict[str, object]) -> dict[str, Path]:
    plt.rcParams["font.family"] = "DejaVu Sans"
    paths: dict[str, Path] = {}

    top = scored[scored["new_atm_candidate"]].head(15).sort_values("placement_score")
    fig, ax = plt.subplots(figsize=(11, 7))
    ax.barh(top["h3_index"], top["placement_score"], color="#0057b8")
    ax.set_title("v2: топ-15 зон для новых банкоматов по ML placement score", fontsize=15, fontweight="bold")
    ax.set_xlabel("placement_score, 0-100")
    ax.grid(axis="x", alpha=0.25)
    for i, value in enumerate(top["placement_score"]):
        ax.text(value + 0.4, i, f"{value:.1f}", va="center", fontsize=9)
    fig.tight_layout()
    paths["top_shortlist"] = CHART_OUT / "v2_top_shortlist.png"
    fig.savefig(paths["top_shortlist"], dpi=180)
    plt.close(fig)

    fig, ax = plt.subplots(figsize=(10, 5.8))
    ax.hist(scored["ml_atm_probability"], bins=35, color="#4d8fd6", edgecolor="white")
    ax.axvline(scored["ml_atm_probability"].quantile(0.8), color="#f05a28", linewidth=2, label="80-й перцентиль")
    ax.set_title("Распределение ML-вероятности ATM-сигнала", fontsize=15, fontweight="bold")
    ax.set_xlabel("ml_atm_probability")
    ax.set_ylabel("Количество H3-зон")
    ax.legend()
    ax.grid(axis="y", alpha=0.25)
    fig.tight_layout()
    paths["ml_probability_distribution"] = CHART_OUT / "v2_ml_probability_distribution.png"
    fig.savefig(paths["ml_probability_distribution"], dpi=180)
    plt.close(fig)

    top_imp = importances.head(15).sort_values("importance")
    fig, ax = plt.subplots(figsize=(10, 6.2))
    ax.barh(top_imp["feature"], top_imp["importance"], color="#2c7fb8")
    ax.set_title("Топ-15 признаков RandomForest", fontsize=15, fontweight="bold")
    ax.set_xlabel("Feature importance")
    ax.grid(axis="x", alpha=0.25)
    fig.tight_layout()
    paths["feature_importance"] = CHART_OUT / "v2_feature_importance.png"
    fig.savefig(paths["feature_importance"], dpi=180)
    plt.close(fig)

    fig, ax = plt.subplots(figsize=(8.5, 5.2))
    labels = ["ROC-AUC", "PR-AUC", "Precision top20%", "Recall top20%"]
    values = [
        model_metrics["roc_auc"],
        model_metrics["average_precision"],
        model_metrics["precision_top20pct"],
        model_metrics["recall_top20pct"],
    ]
    ax.bar(labels, values, color=["#0057b8", "#2c7fb8", "#78b7ff", "#f05a28"])
    ax.set_ylim(0, 1)
    ax.set_title("Качество ML-модели на holdout", fontsize=15, fontweight="bold")
    ax.grid(axis="y", alpha=0.25)
    for i, value in enumerate(values):
        ax.text(i, value + 0.02, f"{value:.3f}", ha="center", fontsize=10)
    fig.tight_layout()
    paths["model_metrics"] = CHART_OUT / "v2_model_metrics.png"
    fig.savefig(paths["model_metrics"], dpi=180)
    plt.close(fig)
    return paths


def h3_feature(row: pd.Series) -> dict[str, object]:
    boundary = h3.cell_to_boundary(row["h3_index"])
    coords = [[lng, lat] for lat, lng in boundary]
    coords.append(coords[0])

    def safe_float(name: str, digits: int = 3):
        value = row.get(name, None)
        if pd.isna(value):
            return None
        return round(float(value), digits)

    props = {
        "h3_index": str(row["h3_index"]),
        "lat": round(float(row["lat"]), 7),
        "lon": round(float(row["lon"]), 7),
        "rank": int(row["priority_rank"]),
        "placement_score": round(float(row["placement_score"]), 2),
        "ml_probability": round(float(row["ml_atm_probability"]), 3),
        "demand_score_v2": round(float(row["demand_score_v2"]), 3),
        "coverage_gap_score": round(float(row["coverage_gap_score"]), 3),
        "tx_frequency_pct": round(float(row.get("tx_frequency_pct", 0)), 1),
        "cell_y": int(row["cell_y"]),
        "recommendation_type": str(row["recommendation_type"]),
        "why_recommended": str(row["why_recommended"]),
        "customers": int(row["unique_customers"]),
        "tx_count": int(row["total_tx_count"]),
        "amount": round(float(row["total_tx_sum"]), 2),
        "avg_ticket": safe_float("avg_ticket_weighted", 2),
        "tx_per_customer": safe_float("tx_per_customer", 2),
        "sum_per_customer": safe_float("sum_per_customer", 2),
        "transaction_volatility": safe_float("transaction_volatility", 2),
        "unique_mcc_count": int(row["unique_mcc_count"]),
        "mcc_entropy": safe_float("mcc_entropy", 3),
        "top_mcc_share": safe_float("top_mcc_share", 3),
        "active_time_buckets": int(row["active_time_buckets"]),
        "time_entropy": safe_float("time_entropy", 3),
        "peak_bucket_share": safe_float("peak_bucket_share", 3),
        "vtb_atm": int(row["vtb_atm"]),
        "vtb_branch": int(row["vtb_branch"]),
        "competitor_atm": int(row["competitor_atm"]),
        "competitor_bank": int(row["competitor_bank"]),
        "mall": int(row["mall"]),
        "transit": int(row["transit"]),
        "education": int(row["education"]),
        "distance_to_nearest_vtb_atm": safe_float("distance_to_nearest_vtb_atm", 3),
        "distance_to_nearest_competitor_atm": safe_float("distance_to_nearest_competitor_atm", 3),
        "metric_formula": "placement_score = 100 × (0.55 × ML probability + 0.25 × demand_score_v2 + 0.20 × coverage_gap_score)",
        "new_atm_candidate": bool(row.get("new_atm_candidate", False)),
        "capacity_candidate": bool(row.get("capacity_candidate", False)),
    }
    return {"type": "Feature", "geometry": {"type": "Polygon", "coordinates": [coords]}, "properties": props}



def metric_popup_html(title: str, label: str, description: str, formula: str, name: str | None = None, detail_label: str = "OSM") -> str:
    name_block = f"<div style='margin-top:6px;color:#53616f;'>{escape(detail_label)}: {escape(name)}</div>" if name else ""
    return f"""
    <div style="font-family:Arial, sans-serif; width:270px; line-height:1.32;">
      <div style="font-size:15px; font-weight:700; color:#1f4e78; margin-bottom:2px;">{escape(title)}</div>
      <div style="font-size:12px; font-weight:700; color:#2f3a45; margin-bottom:6px;">{escape(label)}</div>
      <div style="font-size:12px; color:#2f3a45;">{escape(description)}</div>
      <div style="font-size:11px; color:#6c7785; margin-top:8px; padding-top:6px; border-top:1px solid #d9e2f3;">{escape(formula)}</div>
      {name_block}
    </div>
    """



def add_h3_search_control(fmap: folium.Map, scored: pd.DataFrame) -> None:
    ui_scored = scored.copy()
    ui_scored["tx_frequency_pct"] = ui_scored["total_tx_count"].rank(pct=True).mul(100)
    geojson = {"type": "FeatureCollection", "features": [h3_feature(row) for _, row in ui_scored.iterrows()]}
    map_name = fmap.get_name()
    geojson_json = json.dumps(geojson, ensure_ascii=False)

    html = f"""
    <style>
      :root {{
        --vtb-blue:#003f88;
        --vtb-blue2:#0057b8;
        --vtb-bg:#f6f8fb;
        --vtb-card:rgba(255,255,255,.97);
        --vtb-border:#dbe5f0;
        --vtb-text:#172033;
        --vtb-muted:#687789;
        --vtb-shadow:0 14px 34px rgba(15,33,56,.16);
      }}
      .vtb-topbar {{
        position:absolute; top:14px; left:58px; z-index:1000;
        display:flex; align-items:center; gap:12px;
        background:var(--vtb-card); border:1px solid var(--vtb-border);
        border-radius:16px; box-shadow:var(--vtb-shadow);
        padding:10px 14px; font-family:Arial, sans-serif; color:var(--vtb-text);
        backdrop-filter: blur(8px);
      }}
      .vtb-logo {{
        width:48px; height:30px; border-radius:10px;
        background:linear-gradient(135deg,#003f88,#0078ff); color:#fff;
        display:flex; align-items:center; justify-content:center;
        font-weight:900; font-size:13px; letter-spacing:.7px;
      }}
      .vtb-title {{font-size:13px; font-weight:900; line-height:1.2; white-space:nowrap;}}
      .vtb-subtitle {{font-size:11px; color:var(--vtb-muted); margin-top:2px;}}
      .vtb-switch {{position:relative; display:inline-block; width:48px; height:26px;}}
      .vtb-switch input {{display:none;}}
      .vtb-slider {{position:absolute; cursor:pointer; inset:0; background:#d9e4f1; border-radius:999px; transition:.18s;}}
      .vtb-slider:before {{content:""; position:absolute; width:20px; height:20px; left:3px; top:3px; background:white; border-radius:50%; box-shadow:0 3px 8px rgba(0,0,0,.24); transition:.18s;}}
      .vtb-switch input:checked + .vtb-slider {{background:var(--vtb-blue2);}}
      .vtb-switch input:checked + .vtb-slider:before {{transform:translateX(22px);}}

      .vtb-panel {{
        position:absolute; top:76px; left:58px; z-index:999; width:360px;
        max-height:calc(100vh - 106px); overflow:auto;
        background:var(--vtb-card); border:1px solid var(--vtb-border);
        border-radius:20px; box-shadow:var(--vtb-shadow); padding:16px;
        font-family:Arial, sans-serif; color:var(--vtb-text); backdrop-filter:blur(8px);
      }}
      .vtb-panel.hidden {{display:none;}}
      .vtb-block {{border-bottom:1px solid #edf2f7; padding-bottom:14px; margin-bottom:14px;}}
      .vtb-block:last-child {{border-bottom:0; padding-bottom:0; margin-bottom:0;}}
      .vtb-block-title {{font-size:11px; font-weight:900; text-transform:uppercase; letter-spacing:.05em; color:var(--vtb-blue); margin-bottom:10px;}}
      .vtb-row {{display:flex; gap:8px; align-items:center;}}
      .vtb-input {{flex:1; min-width:0; border:1px solid #c8d5e3; border-radius:12px; padding:10px 11px; font-size:12px; outline:none; background:white; color:var(--vtb-text);}}
      .vtb-input:focus {{border-color:var(--vtb-blue2); box-shadow:0 0 0 3px rgba(0,87,184,.12);}}
      .vtb-btn {{border:0; border-radius:12px; padding:10px 13px; background:var(--vtb-blue2); color:white; cursor:pointer; font-size:12px; font-weight:900;}}
      .vtb-btn:hover {{background:var(--vtb-blue);}}
      .vtb-status {{font-size:11px; color:var(--vtb-muted); line-height:1.35; min-height:15px; margin-top:8px;}}
      .vtb-label-line {{display:flex; align-items:center; justify-content:space-between; gap:10px; font-size:12px; margin:8px 0 7px; color:#2b394b;}}
      .vtb-badge {{display:inline-flex; align-items:center; justify-content:center; min-width:34px; height:22px; border-radius:999px; background:#eef5ff; color:var(--vtb-blue); font-weight:900; font-size:11px; padding:0 7px;}}
      .vtb-range {{width:100%; accent-color:var(--vtb-blue2);}}
      .vtb-select {{width:100%; border:1px solid #c8d5e3; border-radius:12px; padding:9px 10px; background:white; font-size:12px; color:var(--vtb-text); outline:none;}}
      .vtb-hint {{font-size:11px; color:var(--vtb-muted); line-height:1.35; margin-top:8px;}}
      .vtb-segments {{display:grid; grid-template-columns:1fr 1fr; gap:7px;}}
      .vtb-segment {{border:1px solid #d8e4f2; background:#f8fbff; color:#263548; border-radius:12px; padding:9px 8px; font-size:11px; font-weight:800; cursor:pointer; text-align:left;}}
      .vtb-segment.active {{border-color:var(--vtb-blue2); background:#eaf4ff; color:var(--vtb-blue); box-shadow:0 0 0 3px rgba(0,87,184,.08);}}
      .vtb-grid {{display:grid; grid-template-columns:1fr 1fr; gap:8px;}}
      .vtb-card {{background:#f7faff; border:1px solid #e1ebf6; border-radius:14px; padding:10px;}}
      .vtb-card-value {{font-size:16px; font-weight:900; color:var(--vtb-blue);}}
      .vtb-card-label {{font-size:10px; color:var(--vtb-muted); margin-top:3px; line-height:1.25;}}
      .vtb-toplist {{display:flex; flex-direction:column; gap:8px;}}
      .vtb-zone-card {{border:1px solid #e0eaf5; background:#fbfdff; border-radius:14px; padding:9px 10px; cursor:pointer; transition:.15s;}}
      .vtb-zone-card:hover {{border-color:var(--vtb-blue2); box-shadow:0 7px 18px rgba(0,87,184,.12); transform:translateY(-1px);}}
      .vtb-zone-title {{display:flex; justify-content:space-between; gap:8px; font-size:12px; font-weight:900; color:#172033;}}
      .vtb-zone-meta {{margin-top:4px; font-size:11px; color:#687789; line-height:1.35;}}

      .vtb-inspector {{
        position:absolute; top:82px; right:24px; z-index:1002; width:430px;
        max-height:calc(100vh - 115px); overflow:auto; display:none;
        background:var(--vtb-card); border:1px solid var(--vtb-border); border-radius:22px;
        box-shadow:0 18px 44px rgba(15,33,56,.22); padding:18px;
        font-family:Arial, sans-serif; color:var(--vtb-text); backdrop-filter:blur(9px);
      }}
      .vtb-inspector.open {{display:block;}}
      .vtb-inspector-head {{display:flex; align-items:flex-start; justify-content:space-between; gap:12px; margin-bottom:10px;}}
      .vtb-inspector-title {{font-size:18px; line-height:1.2; font-weight:900; color:var(--vtb-blue);}}
      .vtb-inspector-close {{border:0; width:30px; height:30px; border-radius:10px; background:#eef3f9; cursor:pointer; color:#5c6a7b; font-size:18px; line-height:1;}}
      .vtb-inspector-close:hover {{background:#dfeaf6; color:#1b2d44;}}
      .vtb-inspector-badge {{display:inline-block; margin:0 0 12px; padding:6px 10px; border-radius:999px; background:#eef5ff; color:var(--vtb-blue); font-size:12px; font-weight:900;}}
      .vtb-section-title-small {{font-size:12px; font-weight:900; color:#003f88; margin:14px 0 8px; text-transform:uppercase; letter-spacing:.04em;}}
      .vtb-formula {{background:#f5f9ff; border:1px solid #dceafe; border-radius:14px; padding:11px 12px; font-size:12px; line-height:1.45; color:#273548;}}
      .vtb-metric-grid {{display:grid; grid-template-columns:1fr 1fr; gap:8px;}}
      .vtb-metric-card {{background:#fff; border:1px solid #e2ebf5; border-radius:14px; padding:10px;}}
      .vtb-metric-label {{font-size:10px; line-height:1.25; color:#687789; font-weight:800;}}
      .vtb-metric-value {{font-size:15px; font-weight:900; color:#172033; margin-top:5px;}}
      .vtb-report {{font-size:12px; line-height:1.55; color:#35445a; background:#fbfcfe; border:1px solid #e3ecf6; border-radius:14px; padding:11px 12px;}}
      .vtb-report b {{color:#172033;}}

      .leaflet-control-layers {{
        border-radius:18px!important; border:1px solid var(--vtb-border)!important;
        box-shadow:var(--vtb-shadow)!important; font-family:Arial,sans-serif!important;
        font-size:12px!important; overflow:hidden!important; background:var(--vtb-card)!important;
      }}
      .leaflet-control-layers-expanded {{padding:14px 16px!important; background:var(--vtb-card)!important;}}
      .leaflet-control-layers-base, .leaflet-control-layers-separator {{display:none!important;}}
      .leaflet-control-layers-overlays {{margin:0!important;}}
      .vtb-layer-clean-title {{font-size:12px; font-weight:900; color:var(--vtb-blue); margin:0 0 10px; text-transform:uppercase; letter-spacing:.04em;}}
      .vtb-layer-clean-row {{
        display:grid; grid-template-columns:18px 20px 1fr; align-items:center; column-gap:8px;
        font-size:12px; line-height:1.25; margin:9px 0; color:#263548; cursor:pointer; user-select:none;
      }}
      .vtb-layer-clean-row input {{width:15px; height:15px; margin:0; accent-color:var(--vtb-blue2); cursor:pointer;}}
      .vtb-layer-clean-row:hover .vtb-layer-clean-text {{color:var(--vtb-blue);}}
      .vtb-layer-clean-text {{white-space:nowrap;}}
      .vtb-dot {{width:17px; height:17px; border-radius:50%; display:inline-flex; align-items:center; justify-content:center; color:white; font-size:8px; font-weight:900; flex:0 0 auto; box-shadow:0 1px 4px rgba(0,0,0,.22);}}
      .b-blue{{background:#0057b8}} .b-dark{{background:#003f88}} .b-red{{background:#d93025}} .b-orange{{background:#f29900}} .b-purple{{background:#7b4de3}} .b-green{{background:#1e9b4b}} .b-cadet{{background:#2486a6}}
    </style>

    <script>
      window.addEventListener("load", function() {{
        const map = {map_name};
        const h3Data = {geojson_json};
        let h3Layer = null;
        let activeScenario = "high";
        let selectedLayer = null;
        const metricList = [
          ["placement_score", "Итоговый потенциал"],
          ["ml_probability", "ML-вероятность ATM-сигнала"],
          ["demand_score_v2", "Транзакционный спрос"],
          ["coverage_gap_score", "Разрыв покрытия ВТБ"],
          ["tx_frequency_pct", "Частотность операций, перцентиль"],
          ["customers", "Уникальные клиенты"],
          ["tx_count", "Количество операций"],
          ["amount", "Оборот"],
          ["avg_ticket", "Средний чек"],
          ["unique_mcc_count", "Количество MCC"],
          ["mcc_entropy", "Разнообразие MCC"],
          ["time_entropy", "Равномерность по времени"],
          ["vtb_atm", "Банкоматы ВТБ в H3"],
          ["competitor_atm", "Банкоматы конкурентов"],
          ["mall", "Торговые центры"],
          ["education", "Образовательные объекты"],
          ["distance_to_nearest_vtb_atm", "До ближайшего ATM ВТБ, км"]
        ];

        function formatValue(value) {{
          if (value === null || value === undefined || Number.isNaN(value)) return "—";
          if (typeof value === "number") {{
            return value.toLocaleString("ru-RU", {{maximumFractionDigits: 3}});
          }}
          return String(value);
        }}
        function escapeHtml(value) {{
          return String(value ?? "—").replace(/&/g,"&amp;").replace(/</g,"&lt;").replace(/>/g,"&gt;").replace(/\"/g,"&quot;").replace(/'/g,"&#039;");
        }}
        function scoreColor(score) {{
          if (score >= 80) return "#08306b";
          if (score >= 60) return "#2171b5";
          if (score >= 40) return "#6baed6";
          if (score >= 20) return "#c6dbef";
          return "#f7fbff";
        }}
        function currentFilter() {{
          return {{
            minScore: Number(document.querySelector("#vtbMinScore")?.value || 0),
            rankLimit: document.querySelector("#vtbRankLimit")?.value || "all"
          }};
        }}
        function scenarioPass(p) {{
          if (activeScenario === "all") return true;
          if (activeScenario === "high") return Number(p.tx_frequency_pct) >= 80;
          if (activeScenario === "low") return Number(p.tx_frequency_pct) <= 20;
          if (activeScenario === "gap") return Number(p.vtb_atm) === 0 && Number(p.vtb_branch) === 0;
          if (activeScenario === "competitors") return Number(p.competitor_atm) > 0 || Number(p.competitor_bank) > 0;
          if (activeScenario === "new") return Boolean(p.new_atm_candidate);
          return true;
        }}
        function isVisible(p) {{
          const f = currentFilter();
          if (Number(p.placement_score) < f.minScore) return false;
          if (f.rankLimit !== "all" && Number(p.rank) > Number(f.rankLimit)) return false;
          return scenarioPass(p);
        }}
        function visibleFeatures() {{
          return h3Data.features.filter(f => isVisible(f.properties));
        }}
        function styleFeature(feature) {{
          const p = feature.properties;
          const show = isVisible(p);
          return {{
            fillColor: scoreColor(Number(p.placement_score)),
            color: show ? "#274d69" : "#d8e2ec",
            weight: show ? 0.55 : 0.18,
            fillOpacity: show ? 0.62 : 0.045,
            opacity: show ? 0.9 : 0.16
          }};
        }}
        function analyticalText(p) {{
          const vtbPresence = Number(p.vtb_atm) + Number(p.vtb_branch);
          const competitorPresence = Number(p.competitor_atm) + Number(p.competitor_bank);
          const demandText = Number(p.demand_score_v2) >= 0.75
            ? "Транзакционный спрос высокий: зона формирует устойчивый поток операций и может рассматриваться как приоритетная для проверки размещения."
            : Number(p.demand_score_v2) >= 0.45
              ? "Транзакционный спрос умеренный: зона требует дополнительной проверки локации, но уже имеет измеримый клиентский поток."
              : "Транзакционный спрос ниже среднего: зона подходит скорее для мониторинга, чем для первоочередного размещения.";
          const coverageText = vtbPresence === 0
            ? `Внутри H3 не найдено присутствие ВТБ, при этом расстояние до ближайшего ATM ВТБ составляет ${{formatValue(p.distance_to_nearest_vtb_atm)}} км.`
            : `Внутри H3 уже есть присутствие ВТБ: банкоматов ${{formatValue(p.vtb_atm)}}, отделений ${{formatValue(p.vtb_branch)}}.`;
          const competitionText = competitorPresence > 0
            ? `Конкурентная инфраструктура подтверждает рыночный спрос: банкоматов конкурентов ${{formatValue(p.competitor_atm)}}, банков-конкурентов ${{formatValue(p.competitor_bank)}}.`
            : "Прямых конкурентных ATM/банк-точек в H3 не найдено, поэтому спрос лучше дополнительно подтвердить полевым осмотром или внутренними данными.";
          return `
            <b>Аргументация для отчета.</b> Зона имеет итоговый потенциал <b>${{formatValue(p.placement_score)}}</b> при ML-вероятности ATM-сигнала <b>${{formatValue(p.ml_probability)}}</b>. Частотность операций находится на <b>${{formatValue(p.tx_frequency_pct)}}-м перцентиле</b>: зафиксировано <b>${{formatValue(p.tx_count)}}</b> операций, <b>${{formatValue(p.customers)}}</b> уникальных клиентов и оборот <b>${{formatValue(p.amount)}}</b>. ${{demandText}} ${{coverageText}} ${{competitionText}} Разнообразие MCC равно <b>${{formatValue(p.mcc_entropy)}}</b>, равномерность по времени — <b>${{formatValue(p.time_entropy)}}</b>, что помогает оценить устойчивость спроса, а не разовый всплеск.
          `;
        }}
        function inspectorHtml(p) {{
          const metrics = metricList.map(([key, label]) => `
            <div class="vtb-metric-card">
              <div class="vtb-metric-label">${{escapeHtml(label)}}</div>
              <div class="vtb-metric-value">${{formatValue(p[key])}}</div>
            </div>`).join("");
          return `
            <div class="vtb-inspector-head">
              <div class="vtb-inspector-title">H3-зона ${{escapeHtml(p.h3_index)}}</div>
              <button id="vtbInspectorClose" class="vtb-inspector-close" type="button">×</button>
            </div>
            <div class="vtb-inspector-badge">Ранг #${{formatValue(p.rank)}} · ${{escapeHtml(p.recommendation_type)}}</div>
            <div class="vtb-section-title-small">Формула скоринга</div>
            <div class="vtb-formula">${{escapeHtml(p.metric_formula)}}<br><b>Расчет:</b> 100 × (0.55 × ${{formatValue(p.ml_probability)}} + 0.25 × ${{formatValue(p.demand_score_v2)}} + 0.20 × ${{formatValue(p.coverage_gap_score)}}) = <b>${{formatValue(p.placement_score)}}</b></div>
            <div class="vtb-section-title-small">Почему рекомендовано</div>
            <div class="vtb-report">${{analyticalText(p)}}<br><br><b>Краткая причина модели:</b> ${{escapeHtml(p.why_recommended)}}.</div>
            <div class="vtb-section-title-small">Данные зоны</div>
            <div class="vtb-metric-grid">${{metrics}}</div>
          `;
        }}
        function openInspector(p) {{
          const panel = document.querySelector("#vtbInspector");
          panel.innerHTML = inspectorHtml(p);
          panel.classList.add("open");
          document.querySelector("#vtbInspectorClose").addEventListener("click", closeInspector);
        }}
        function closeInspector() {{
          const panel = document.querySelector("#vtbInspector");
          panel.classList.remove("open");
          panel.innerHTML = "";
          if (selectedLayer) {{ map.removeLayer(selectedLayer); selectedLayer = null; }}
        }}
        function focusFeature(feature, openCard=true) {{
          if (selectedLayer) map.removeLayer(selectedLayer);
          selectedLayer = L.geoJSON(feature, {{style: {{color:"#ff7a00", weight:4, fillColor:"#ffd166", fillOpacity:.20}}}}).addTo(map);
          map.fitBounds(selectedLayer.getBounds(), {{maxZoom:13, padding:[28,28]}});
          if (openCard) openInspector(feature.properties);
        }}
        function bindFeature(feature, layer) {{
          layer.on("click", () => {{ if (isVisible(feature.properties)) focusFeature(feature, true); }});
        }}
        function redraw() {{
          if (h3Layer) map.removeLayer(h3Layer);
          h3Layer = L.geoJSON(h3Data, {{style: styleFeature, onEachFeature: bindFeature}}).addTo(map);
          updateDashboard();
          updateTopList();
        }}
        function setLabels() {{
          const minScore = document.querySelector("#vtbMinScore");
          const minScoreValue = document.querySelector("#vtbMinScoreValue");
          if (minScore && minScoreValue) minScoreValue.textContent = minScore.value;
        }}
        function searchH3() {{
          const input = document.querySelector("#vtbTerritorySearch");
          const status = document.querySelector("#vtbSearchStatus");
          const value = input.value.trim().toLowerCase();
          if (!value) {{ status.textContent = "Введите H3-индекс."; return; }}
          const feature = h3Data.features.find(f => String(f.properties.h3_index).toLowerCase() === value);
          if (!feature) {{ status.textContent = "Зона не найдена в расчетной витрине."; return; }}
          status.textContent = `Найдено: ранг #${{feature.properties.rank}}, score ${{formatValue(feature.properties.placement_score)}}`;
          focusFeature(feature, true);
        }}
        function updateDashboard() {{
          const arr = visibleFeatures().map(f => f.properties);
          const count = arr.length;
          const avgScore = count ? arr.reduce((s,p) => s + Number(p.placement_score), 0) / count : 0;
          const highCount = arr.filter(p => Number(p.placement_score) >= 60).length;
          const noVtbCount = arr.filter(p => Number(p.vtb_atm) + Number(p.vtb_branch) === 0).length;
          const competitorCount = arr.filter(p => Number(p.competitor_atm) + Number(p.competitor_bank) > 0).length;
          const set = (id, value) => {{ const el = document.querySelector(id); if (el) el.textContent = value; }};
          set("#vtbVisibleCount", formatValue(count));
          set("#vtbAvgScore", formatValue(Number(avgScore.toFixed(1))));
          set("#vtbHighCount", formatValue(highCount));
          set("#vtbNoVtbCount", formatValue(noVtbCount));
          set("#vtbCompetitorCount", formatValue(competitorCount));
        }}
        function updateTopList() {{
          const box = document.querySelector("#vtbTopList");
          if (!box) return;
          const top = visibleFeatures()
            .sort((a,b) => Number(a.properties.rank) - Number(b.properties.rank))
            .slice(0,5);
          if (!top.length) {{ box.innerHTML = `<div class="vtb-status">Нет зон под выбранные фильтры.</div>`; return; }}
          box.innerHTML = top.map(f => `
            <div class="vtb-zone-card" data-h3="${{escapeHtml(f.properties.h3_index)}}">
              <div class="vtb-zone-title"><span>#${{formatValue(f.properties.rank)}} · ${{formatValue(f.properties.placement_score)}}</span><span>ML ${{formatValue(f.properties.ml_probability)}}</span></div>
              <div class="vtb-zone-meta">${{escapeHtml(f.properties.recommendation_type)}}<br>${{formatValue(f.properties.customers)}} клиентов · ${{formatValue(f.properties.tx_count)}} операций · ${{formatValue(f.properties.tx_frequency_pct)}} перцентиль</div>
            </div>`).join("");
          box.querySelectorAll(".vtb-zone-card").forEach(card => {{
            card.addEventListener("click", () => {{
              const h3 = card.getAttribute("data-h3");
              const feature = h3Data.features.find(f => String(f.properties.h3_index) === h3);
              if (feature) focusFeature(feature, true);
            }});
          }});
        }}
        function enhanceLayerControl() {{
          const control = document.querySelector(".leaflet-control-layers");
          const list = document.querySelector(".leaflet-control-layers-list");
          const overlays = document.querySelector(".leaflet-control-layers-overlays");
          if (!control || !list || !overlays || overlays.querySelector(".vtb-layer-clean-title")) return;

          control.classList.add("vtb-layer-control-clean");

          const sourceLabels = Array.from(overlays.querySelectorAll("label"));
          const rows = [
            {{ match: "Шорт-лист", text: "Шорт-лист новых точек", icon: "★", cls: "b-blue" }},
            {{ match: "Банкоматы ВТБ", text: "Банкоматы ВТБ", icon: "▣", cls: "b-blue" }},
            {{ match: "Отделения ВТБ", text: "Отделения ВТБ", icon: "◆", cls: "b-dark" }},
            {{ match: "Банкоматы конкурентов", text: "Банкоматы конкурентов", icon: "▣", cls: "b-red" }},
            {{ match: "Банки-конкуренты", text: "Банки-конкуренты", icon: "◆", cls: "b-orange" }},
            {{ match: "ТЦ", text: "Торговые центры", icon: "ТЦ", cls: "b-purple" }},
            {{ match: "Транспортные", text: "Транспорт", icon: "М", cls: "b-green" }},
            {{ match: "Образовательные", text: "Образование", icon: "У", cls: "b-cadet" }}
          ];

          const cleanBox = document.createElement("div");
          cleanBox.innerHTML = `<div class="vtb-layer-clean-title">Обозначения слоев</div>`;

          rows.forEach(item => {{
            const oldLabel = sourceLabels.find(label => label.textContent.includes(item.match));
            if (!oldLabel) return;
            const input = oldLabel.querySelector("input");
            if (!input) return;

            const row = document.createElement("label");
            row.className = "vtb-layer-clean-row";

            const dot = document.createElement("span");
            dot.className = `vtb-dot ${{item.cls}}`;
            dot.textContent = item.icon;

            const text = document.createElement("span");
            text.className = "vtb-layer-clean-text";
            text.textContent = item.text;

            row.appendChild(input);
            row.appendChild(dot);
            row.appendChild(text);
            cleanBox.appendChild(row);
          }});

          overlays.replaceChildren(cleanBox);
        }}
        function makeUI() {{
          const shell = document.createElement("div");
          shell.innerHTML = `
            <div class="vtb-topbar">
              <div class="vtb-logo">ВТБ</div>
              <div><div class="vtb-title">Геоаналитика банкоматов · v2 ML</div><div class="vtb-subtitle">Интерактивная карта приоритетов размещения</div></div>
              <label class="vtb-switch"><input id="vtbPanelToggle" type="checkbox" checked><span class="vtb-slider"></span></label>
            </div>
            <div id="vtbPanel" class="vtb-panel">
              <div class="vtb-block">
                <div class="vtb-block-title">Поиск территории</div>
                <div class="vtb-row"><input id="vtbTerritorySearch" class="vtb-input" type="text" placeholder="Например: 8911aa6257bffff"><button id="vtbSearchBtn" class="vtb-btn" type="button">Найти</button></div>
                <div id="vtbSearchStatus" class="vtb-status">Введите H3 и нажмите Enter или «Найти».</div>
              </div>
              <div class="vtb-block">
                <div class="vtb-block-title">Сценарий анализа</div>
                <div class="vtb-segments">
                  <button class="vtb-segment" data-scenario="all" type="button">Все H3-зоны</button>
                  <button class="vtb-segment active" data-scenario="high" type="button">Высокочастотные</button>
                  <button class="vtb-segment" data-scenario="low" type="button">Низкочастотные</button>
                  <button class="vtb-segment" data-scenario="gap" type="button">Без ВТБ внутри H3</button>
                  <button class="vtb-segment" data-scenario="competitors" type="button">Есть конкуренты</button>
                  <button class="vtb-segment" data-scenario="new" type="button">Новые точки</button>
                </div>
                <div class="vtb-hint">Сценарии быстро переключают рабочую выборку: верхние 20% по частотности, нижние 20%, зоны с gap покрытия, конкурентные кластеры и кандидаты на новые точки.</div>
              </div>
              <div class="vtb-block">
                <div class="vtb-block-title">Приоритет и масштаб выборки</div>
                <div class="vtb-label-line"><span>Минимальный placement score</span><span id="vtbMinScoreValue" class="vtb-badge">0</span></div>
                <input id="vtbMinScore" class="vtb-range" type="range" min="0" max="100" value="0" step="1">
                <div style="height:10px"></div>
                <select id="vtbRankLimit" class="vtb-select">
                  <option value="30">Только топ-30 по приоритету</option>
                  <option value="100">Только топ-100 по приоритету</option>
                  <option value="300">Только топ-300 по приоритету</option>
                  <option value="all" selected>Все зоны выбранного сценария</option>
                </select>
              </div>
              <div class="vtb-block">
                <div class="vtb-block-title">Сводка по текущей выборке</div>
                <div class="vtb-grid">
                  <div class="vtb-card"><div id="vtbVisibleCount" class="vtb-card-value">—</div><div class="vtb-card-label">зон на карте</div></div>
                  <div class="vtb-card"><div id="vtbAvgScore" class="vtb-card-value">—</div><div class="vtb-card-label">средний score</div></div>
                  <div class="vtb-card"><div id="vtbHighCount" class="vtb-card-value">—</div><div class="vtb-card-label">зон со score ≥ 60</div></div>
                  <div class="vtb-card"><div id="vtbNoVtbCount" class="vtb-card-value">—</div><div class="vtb-card-label">зон без ВТБ внутри H3</div></div>
                  <div class="vtb-card"><div id="vtbCompetitorCount" class="vtb-card-value">—</div><div class="vtb-card-label">зон с конкурентами</div></div>
                  <div class="vtb-card"><div class="vtb-card-value">20%</div><div class="vtb-card-label">граница high/low частотности</div></div>
                </div>
              </div>
              <div class="vtb-block">
                <div class="vtb-block-title">Топ зон в выбранной выборке</div>
                <div id="vtbTopList" class="vtb-toplist"></div>
              </div>
            </div>
            <div id="vtbInspector" class="vtb-inspector"></div>`;
          document.body.appendChild(shell);
          document.querySelector("#vtbPanelToggle").addEventListener("change", e => document.querySelector("#vtbPanel").classList.toggle("hidden", !e.target.checked));
          document.querySelector("#vtbSearchBtn").addEventListener("click", searchH3);
          document.querySelector("#vtbTerritorySearch").addEventListener("keydown", e => {{ if (e.key === "Enter") searchH3(); }});
          document.querySelectorAll(".vtb-segment").forEach(btn => btn.addEventListener("click", () => {{
            document.querySelectorAll(".vtb-segment").forEach(b => b.classList.remove("active"));
            btn.classList.add("active");
            activeScenario = btn.getAttribute("data-scenario");
            redraw();
          }}));
          ["#vtbMinScore", "#vtbRankLimit"].forEach(sel => document.querySelector(sel).addEventListener("input", () => {{ setLabels(); redraw(); }}));
          document.querySelectorAll(".vtb-panel,.vtb-topbar,.vtb-inspector").forEach(el => {{ L.DomEvent.disableClickPropagation(el); L.DomEvent.disableScrollPropagation(el); }});
          document.addEventListener("keydown", e => {{ if (e.key === "Escape") closeInspector(); }});
          setLabels();
        }}
        makeUI();
        redraw();
        enhanceLayerControl();
        setTimeout(enhanceLayerControl, 250);
      }});
    </script>
    """
    fmap.get_root().html.add_child(Element(html))



def save_map(scored: pd.DataFrame, poi: pd.DataFrame) -> Path:
    center = [float(scored["lat"].mean()), float(scored["lon"].mean())]
    fmap = folium.Map(
    location=center,
    zoom_start=10,
    tiles=None,
    control_scale=True,
    prefer_canvas=True,
    attribution_control=False,)
    folium.TileLayer("CartoDB positron", name="Картографическая подложка", control=True).add_to(fmap)

    colormap = LinearColormap(["#f7fbff", "#c6dbef", "#6baed6", "#2171b5", "#08306b"], vmin=0, vmax=100)
    colormap.caption = "v2 ML placement_score"

    # H3-гексагоны рисуются кастомным JS-слоем: без hover, с карточкой только по клику и с фильтрами слева.
    add_h3_search_control(fmap, scored)

    top_layer = folium.FeatureGroup(name="🔵 ★ Шорт-лист новых точек v2", show=True)
    for _, row in scored[scored["new_atm_candidate"]].head(30).iterrows():
        popup = metric_popup_html(
            "Кандидат на новую точку",
            "Приоритетная зона для размещения банкомата",
            f"Ранг #{int(row['priority_rank'])}, placement_score {row['placement_score']:.1f}, ML probability {row['ml_atm_probability']:.3f}.",
            "placement_score = 100 × (0.55 × ML + 0.25 × demand + 0.20 × coverage_gap)",
            row["why_recommended"],
            "Причина",
        )
        folium.Marker(
            [row["lat"], row["lon"]],
            icon=folium.Icon(color="blue", icon="star"),
            popup=folium.Popup(popup, max_width=340),
        ).add_to(top_layer)
    top_layer.add_to(fmap)

    poi_styles = {
        "vtb_atm": {"color": "blue", "icon": "credit-card", "name": "🔵 ▣ Банкоматы ВТБ"},
        "vtb_branch": {"color": "darkblue", "icon": "briefcase", "name": "🔷 ◆ Отделения ВТБ"},
        "competitor_atm": {"color": "red", "icon": "credit-card", "name": "🔴 ▣ Банкоматы конкурентов"},
        "competitor_bank": {"color": "orange", "icon": "briefcase", "name": "🟠 ◆ Банки-конкуренты"},
        "mall": {"color": "purple", "icon": "shopping-cart", "name": "🟣 ТЦ и торговые объекты"},
        "transit": {"color": "green", "icon": "road", "name": "🟢 Транспортные объекты"},
        "education": {"color": "cadetblue", "icon": "education", "name": "🔵 У Образовательные объекты"},
    }

    for category, style in poi_styles.items():
        subset = poi[(poi["category"] == category) & (poi["in_case_cells"])]
        if subset.empty:
            continue
        cluster = MarkerCluster(name=style["name"], show=False)
        title, label, description, formula = POI_CATEGORY_HELP[category]
        for _, row in subset.iterrows():
            folium.Marker(
                [row["lat"], row["lon"]],
                icon=folium.Icon(color=style["color"], icon=style["icon"], prefix="fa"),
                popup=folium.Popup(
                    metric_popup_html(
                        title=label,
                        label="OSM-объект",
                        description=description,
                        formula=formula,
                        name=row["name"],
                        detail_label="Название",
                    ),
                    max_width=320,
                ),
            ).add_to(cluster)
        cluster.add_to(fmap)

    colormap.add_to(fmap)
    folium.LayerControl(collapsed=False, position="topright").add_to(fmap)
    path = MAP_OUT / "v2_atm_ml_potential_map.html"
    fmap.save(path)
    return path


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.25), Inches(0.58))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = VTB_BLUE
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.9), Inches(11.9), Inches(0.4))
        p2 = sub.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = MUTED


def add_bullets(slide, bullets: list[str], x: float, y: float, w: float, h: float, size: int = 15) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(7)


def add_metric_card(slide, x: float, y: float, w: float, h: float, value: str, label: str) -> None:
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(232, 241, 255)
    shape.line.color.rgb = RGBColor(190, 210, 235)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = value
    p.font.bold = True
    p.font.size = Pt(21)
    p.font.color.rgb = VTB_BLUE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(10)
    p2.font.color.rgb = MUTED
    p2.alignment = PP_ALIGN.CENTER


def add_table(slide, df: pd.DataFrame, x: float, y: float, w: float, h: float) -> None:
    table = slide.shapes.add_table(len(df) + 1, len(df.columns), Inches(x), Inches(y), Inches(w), Inches(h)).table
    for c, col in enumerate(df.columns):
        cell = table.cell(0, c)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = VTB_BLUE
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(8.5)
    for r, (_, row) in enumerate(df.iterrows(), start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            cell.text_frame.paragraphs[0].font.color.rgb = TEXT


def save_presentation(scored: pd.DataFrame, diagnostics: dict[str, object], model_metrics: dict[str, object], charts: dict[str, Path], artifacts: dict[str, Path], map_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_title(slide, "v2: ML-модель потенциала банкоматных зон", "Расширенный parquet-датасет позволяет обучить H3-level классификатор ATM-сигнала")
    add_bullets(
        slide,
        [
            "В отличие от Excel v1, в v2 есть отрицательные H3-зоны: можно обучить модель cell_y.",
            "Модель предсказывает вероятность ATM-сигнала по H3, а бизнес-скор добавляет спрос и gap покрытия ВТБ.",
            "Результат: ML-шорт-лист новых зон, отдельный список усиления сети, карта и воспроизводимые CSV.",
        ],
        0.75,
        1.55,
        7.0,
        3.2,
        17,
    )
    add_metric_card(slide, 8.15, 1.35, 1.55, 1.0, f"{diagnostics['h3_cells_data']}", "H3-зон")
    add_metric_card(slide, 9.9, 1.35, 1.55, 1.0, f"{diagnostics['h3_positive_cells_in_data']}", "positive cells")
    add_metric_card(slide, 11.65, 1.35, 1.55, 1.0, f"{diagnostics['h3_negative_cells_in_data']}", "negative cells")
    add_metric_card(slide, 8.15, 2.65, 1.55, 1.0, f"{model_metrics['roc_auc']:.3f}", "ROC-AUC")
    add_metric_card(slide, 9.9, 2.65, 1.55, 1.0, f"{model_metrics['average_precision']:.3f}", "PR-AUC")
    add_metric_card(slide, 11.65, 2.65, 1.55, 1.0, f"{model_metrics['precision_top20pct']:.3f}", "precision top20%")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Данные и target", "v2 дает корректную постановку ML на уровне H3-зоны")
    add_bullets(
        slide,
        [
            f"data.parquet: {diagnostics['data_rows']:,} строк; target.parquet: {diagnostics['target_rows_unique']:,} уникальных пар.".replace(",", " "),
            f"H3-зон в data: {diagnostics['h3_cells_data']}; positive target cells: {diagnostics['h3_positive_cells_in_data']}; negative cells: {diagnostics['h3_negative_cells_in_data']}.",
            "Цель модели: cell_y = 1, если H3 есть в target, иначе 0.",
            "Часть target-пар отсутствует в data как транзакционная пара; поэтому ML строится на уровне зоны, а не пары customer-cell.",
        ],
        0.75,
        1.35,
        6.2,
        4.8,
        15,
    )
    slide.shapes.add_picture(str(charts["ml_probability_distribution"]), Inches(7.1), Inches(1.35), width=Inches(5.65))

    slide = prs.slides.add_slide(blank)
    add_title(slide, "ML-подход и скоринг", "Модель предсказывает ATM-сигнал, бизнес-слой переводит его в приоритет размещения")
    add_bullets(
        slide,
        [
            "Модель: RandomForestClassifier с class_weight='balanced_subsample'.",
            "Признаки: транзакционный спрос, чековое поведение, MCC-вектор, временная стабильность, OSM-инфраструктура и расстояния.",
            "placement_score = 100 * (0.55 * ML probability + 0.25 * demand_score_v2 + 0.20 * coverage_gap_score).",
            "Фильтр новых точек: нет найденного ВТБ-присутствия в OSM и нет аномалии высокого чека при малой базе.",
        ],
        0.75,
        1.35,
        5.85,
        5.1,
        14,
    )
    slide.shapes.add_picture(str(charts["model_metrics"]), Inches(6.85), Inches(1.45), width=Inches(5.6))

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Что влияет на модель", "RandomForest feature importance")
    slide.shapes.add_picture(str(charts["feature_importance"]), Inches(0.75), Inches(1.2), width=Inches(6.3))
    add_bullets(
        slide,
        [
            "Feature importance используется как sanity check: модель должна опираться на спрос, инфраструктуру и поведенческие признаки.",
            "Окончательная бизнес-рекомендация не равна чистой ML-вероятности: добавлен coverage gap и исключение зон с найденным ВТБ-присутствием.",
        ],
        7.3,
        1.45,
        5.1,
        4.6,
        15,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Шорт-лист новых банкоматов", "Топ-10 из v2_shortlist_new_atm.csv")
    top = scored[scored["new_atm_candidate"]].head(10)
    table = pd.DataFrame(
        {
            "Ранг": top["priority_rank"],
            "H3": top["h3_index"],
            "Score": top["placement_score"].round(1),
            "ML": top["ml_atm_probability"].round(3),
            "Target": top["cell_y"],
            "Клиенты": top["unique_customers"],
            "Тип": top["recommendation_type"].str.slice(0, 28),
        }
    )
    add_table(slide, table, 0.35, 1.25, 12.7, 4.45)
    add_bullets(slide, ["Полный список: output/data/v2_shortlist_new_atm.csv"], 0.6, 6.05, 11.7, 0.5, 13)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Карта и интерфейс", "HTML-карта с поиском H3, разными hover/click данными и OSM-слоями")
    slide.shapes.add_picture(str(charts["top_shortlist"]), Inches(0.65), Inches(1.2), width=Inches(5.95))
    add_bullets(
        slide,
        [
            f"Карта: {map_path.relative_to(ROOT).as_posix()}",
            "Наведение на H3 показывает операционные данные зоны.",
            "Клик по H3 показывает ML-формулу и расшифровку метрик.",
            "Поиск по H3 быстро центрирует карту и подсвечивает зону.",
            "OSM-слои доступны отдельными переключателями.",
        ],
        6.95,
        1.4,
        5.65,
        4.8,
        15,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Внедрение и риски", "Что можно делать дальше")
    add_bullets(
        slide,
        [
            "Внедрение: регулярный пересчет parquet -> признаки H3 -> ML probability -> placement score -> карта/очередь обследования.",
            "Калибровка: заменить proxy target на фактическую эффективность банкоматов ВТБ, если появятся операции/затраты/инкассация.",
            "Риск target: target показывает факт ATM-активности, но не финансовый эффект установки нового банкомата.",
            "Риск OSM: открытые данные неполны; для промышленного контура нужен внутренний справочник ВТБ и проверенный слой конкурентов.",
            "H3-зона - не точный адрес; финальный выбор требует изохрон, аренды, безопасности и техограничений площадки.",
        ],
        0.85,
        1.35,
        11.6,
        5.5,
        15,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Материалы v2")
    lines = [
        f"Все зоны: {artifacts['zone_scores'].relative_to(ROOT).as_posix()}",
        f"Шорт-лист: {artifacts['shortlist'].relative_to(ROOT).as_posix()}",
        f"Feature importance: {artifacts['importances'].relative_to(ROOT).as_posix()}",
        f"Полная витрина признаков: {artifacts['full_features'].relative_to(ROOT).as_posix()}",
        f"Карта: {map_path.relative_to(ROOT).as_posix()}",
        f"Диагностика и метрики модели: {artifacts['diagnostics'].relative_to(ROOT).as_posix()}",
        "Код: scripts/analyze_case_v2_ml.py",
    ]
    add_bullets(slide, lines, 0.85, 1.35, 11.6, 4.8, 15)

    path = PPTX_OUT / "vtb_atm_geoanalytics_solution_v2_ml.pptx"
    try:
        prs.save(path)
    except PermissionError:
        suffix = datetime.now().strftime("%Y%m%d_%H%M%S")
        path = PPTX_OUT / f"vtb_atm_geoanalytics_solution_v2_ml_{suffix}.pptx"
        prs.save(path)
    return path


def save_summary(scored: pd.DataFrame, diagnostics: dict[str, object], model_metrics: dict[str, object], artifacts: dict[str, Path], map_path: Path, pptx_path: Path) -> Path:
    top = scored[scored["new_atm_candidate"]].head(10)
    lines = [
        "# v2 ML-решение кейса",
        "",
        "## Что изменилось относительно первой версии",
        "",
        "- В v2 есть отрицательные H3-зоны: 6500 зон без target-сигнала против 1654 зон с target-сигналом.",
        "- Поэтому ML-модель теперь методологически допустима на уровне H3-зоны.",
        "- Цель модели: `cell_y = 1`, если `h3_index` есть в `target.parquet`, иначе `0`.",
        "",
        "## Качество модели",
        "",
        f"- ROC-AUC: {model_metrics['roc_auc']:.3f}",
        f"- PR-AUC: {model_metrics['average_precision']:.3f}",
        f"- Precision top 20%: {model_metrics['precision_top20pct']:.3f}",
        f"- Recall top 20%: {model_metrics['recall_top20pct']:.3f}",
        "",
        "## Итоговый score",
        "",
        "`placement_score = 100 * (0.55 * ml_atm_probability + 0.25 * demand_score_v2 + 0.20 * coverage_gap_score)`",
        "",
        "## Топ-10 новых зон",
        "",
    ]
    for _, row in top.iterrows():
        lines.append(
            f"- {int(row['priority_rank'])}. `{row['h3_index']}` - score {row['placement_score']:.1f}, "
            f"ML {row['ml_atm_probability']:.3f}, target {int(row['cell_y'])}; {row['why_recommended']}."
        )
    lines.extend(
        [
            "",
            "## Артефакты",
            "",
            f"- Презентация: `{pptx_path.relative_to(ROOT).as_posix()}`",
            f"- Карта: `{map_path.relative_to(ROOT).as_posix()}`",
            f"- Все зоны: `{artifacts['zone_scores'].relative_to(ROOT).as_posix()}`",
            f"- Полная витрина признаков: `{artifacts['full_features'].relative_to(ROOT).as_posix()}`",
            f"- Шорт-лист: `{artifacts['shortlist'].relative_to(ROOT).as_posix()}`",
            f"- Диагностика: `{artifacts['diagnostics'].relative_to(ROOT).as_posix()}`",
            "",
            "## Ограничения",
            "",
            "- Модель предсказывает наличие target-сигнала в H3, а не прямую окупаемость банкомата.",
            "- Target-пары частично отсутствуют в data как customer-cell пары, поэтому задача поставлена на уровне H3, а не пары клиент-зона.",
            "- OSM-слои используются как открытый proxy; для промышленного решения нужен внутренний справочник точек ВТБ и конкурентов.",
        ]
    )
    path = OUT / "v2_solution_summary.md"
    path.write_text("\n".join(lines), encoding="utf-8")
    return path


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="v2 ML solution for VTB ATM placement case")
    parser.add_argument("--refresh-osm", action="store_true", help="Fetch OSM again instead of using output/osm cache")
    parser.add_argument("--no-osm", action="store_true", help="Skip OSM features")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    ensure_dirs()
    data, target = load_data()
    temp_cell = (
        data[["h3_index"]]
        .drop_duplicates()
        .assign(
            lat=lambda df: df["h3_index"].map(lambda idx: h3.cell_to_latlng(idx)[0]),
            lon=lambda df: df["h3_index"].map(lambda idx: h3.cell_to_latlng(idx)[1]),
        )
    )
    raw_osm, osm_meta = fetch_osm(temp_cell, refresh=args.refresh_osm, no_osm=args.no_osm)
    poi = osm_to_poi(raw_osm, set(temp_cell["h3_index"]))
    cell, diagnostics = build_cell_features(data, target, poi)
    _, scored, model_metrics, importances = train_model(cell)
    scored = add_business_scores(scored)
    artifacts = export_artifacts(scored, poi, diagnostics, model_metrics, importances, osm_meta)
    charts = save_charts(scored, importances, model_metrics)
    map_path = save_map(scored, poi)
    pptx_path = save_presentation(scored, diagnostics, model_metrics, charts, artifacts, map_path)
    summary_path = save_summary(scored, diagnostics, model_metrics, artifacts, map_path, pptx_path)
    print("Done v2 ML")
    print(f"Summary: {summary_path}")
    print(f"Presentation: {pptx_path}")
    print(f"Map: {map_path}")
    print(f"Shortlist: {artifacts['shortlist']}")
    print(f"ROC-AUC: {model_metrics['roc_auc']:.4f}")
    print(scored[scored["new_atm_candidate"]].head(10)[["priority_rank", "h3_index", "placement_score", "ml_atm_probability", "cell_y", "recommendation_type"]].to_string(index=False))


if __name__ == "__main__":
    main()
